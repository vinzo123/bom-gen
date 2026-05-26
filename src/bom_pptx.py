from __future__ import annotations

import json
import re
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

try:
    from pptx.enum.dml import MSO_ARROWHEAD
except ImportError:  # pragma: no cover
    MSO_ARROWHEAD = None


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
GREEN = RGBColor(142, 209, 73)
BLACK = RGBColor(12, 12, 12)
INK = RGBColor(39, 52, 64)
MUTED = RGBColor(94, 108, 119)
BLUE = RGBColor(47, 131, 209)
TEAL = RGBColor(24, 169, 155)
RED = RGBColor(173, 98, 102)
BROWN = RGBColor(141, 130, 111)
PANEL = RGBColor(221, 230, 236)
WHITE = RGBColor(255, 255, 255)


def safe_name(value: str, fallback: str = "bom") -> str:
    cleaned = re.sub(r"[^A-Za-z0-9]+", "_", value).strip("_").lower()
    return cleaned or fallback


def set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color: RGBColor = INK, width: float = 1) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_text(slide, x, y, w, h, text: str, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = "Arial"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, bullets: list[str], size=20, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.14)
    frame.margin_right = Inches(0.12)
    frame.margin_top = Inches(0.08)
    frame.margin_bottom = Inches(0.08)

    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"- {bullet}"
        paragraph.space_after = Pt(8)
        paragraph.font.name = "Arial"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
    return box


def add_common_frame(slide, title: str) -> None:
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(248, 248, 246)

    top_black = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.34))
    set_fill(top_black, BLACK)
    top_black.line.fill.background()

    top_green = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.3), Inches(0.09))
    set_fill(top_green, GREEN)
    top_green.line.fill.background()

    bottom_green = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.8), Inches(8.0), Inches(0.7))
    set_fill(bottom_green, GREEN)
    bottom_green.line.fill.background()

    bottom_black = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.15), Inches(6.8), Inches(4.18), Inches(0.7))
    set_fill(bottom_black, BLACK)
    bottom_black.line.fill.background()

    add_text(slide, 0.85, 0.64, 11.65, 0.35, title.upper(), size=24, bold=True, color=RGBColor(0, 0, 0), align=PP_ALIGN.CENTER)
    add_text(slide, 0.75, 6.96, 7.1, 0.22, "ECOSOFT ZOLUTIONS PVT LTD", size=17, bold=True, color=RGBColor(0, 0, 0))
    add_text(slide, 0.75, 7.19, 7.1, 0.16, "Formerly Known as HNS INFOTECH PVT LTD", size=10, bold=True, color=RGBColor(0, 0, 0))


def add_content_border(slide) -> None:
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.08), Inches(11.95), Inches(5.38))
    border.fill.background()
    set_line(border, RGBColor(0, 0, 0), 1.1)


def add_table(slide, x, y, w, h, headers: list[str], rows: list[list[str]], font_size=13):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for col_index, header in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header
        set_fill(cell, GREEN)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(font_size)
            paragraph.font.name = "Arial"
            paragraph.font.color.rgb = RGBColor(0, 0, 0)

    for row_index, row in enumerate(rows, start=1):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Arial"
                paragraph.font.color.rgb = INK
    return table_shape


def add_arrow(slide, x1, y1, x2, y2, color=INK, width=2.5) -> None:
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    if MSO_ARROWHEAD is not None:
        connector.line.end_arrowhead = MSO_ARROWHEAD.TRIANGLE


def add_box(slide, x, y, w, h, text, fill=PANEL, line=INK, size=13, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill)
    set_line(shape, line, 1)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.name = "Arial"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = INK
    return shape


def slide_title(prs: Presentation, content: dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_common_frame(slide, "")
    project = content["project"]
    title = f'{project["project_code"]}_{project["solution_name"]}'
    add_text(slide, 1.15, 1.55, 11.0, 0.7, title, size=31, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(slide, 1.15, 2.45, 11.0, 0.35, "For", size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(slide, 1.15, 2.95, 11.0, 0.85, project["customer_name"], size=26, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(slide, 1.15, 4.15, 11.0, 0.35, "Prepared by Ecosoft Zolutions Pvt Ltd", size=19, bold=True, color=MUTED, align=PP_ALIGN.CENTER)


def slide_text(prs: Presentation, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_common_frame(slide, title)
    add_content_border(slide)
    add_bullets(slide, 1.05, 1.55, 11.15, 4.55, bullets, size=20)


def slide_technology(prs: Presentation, content: dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_common_frame(slide, "TECHNOLOGY")
    add_content_border(slide)
    technologies = content["technology"]
    for index, item in enumerate(technologies[:6]):
        y = 1.55 + (index * 0.68)
        add_box(slide, 1.2, y, 10.85, 0.45, item, fill=RGBColor(238, 244, 238), line=GREEN, size=16, bold=True)


def slide_process_flow(prs: Presentation, content: dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_common_frame(slide, "PROCESS FLOW")
    add_content_border(slide)
    steps = content["process_steps"][:6]
    x_positions = [1.1, 3.05, 5.0, 6.95, 8.9, 10.85]
    y = 2.45
    for index, step in enumerate(steps):
        add_box(slide, x_positions[index], y, 1.45, 0.75, step, fill=RGBColor(231, 240, 247), line=BLUE, size=11, bold=True)
        if index < len(steps) - 1:
            add_arrow(slide, x_positions[index] + 1.45, y + 0.38, x_positions[index + 1], y + 0.38, color=BLUE)

    add_text(slide, 1.05, 4.15, 3.5, 0.25, "Mermaid source generated:", size=12, bold=True, color=MUTED)
    add_text(slide, 1.05, 4.45, 10.9, 1.15, content["process_mermaid"], size=10, color=RGBColor(60, 60, 60))


def slide_solution_illustration(prs: Presentation, content: dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_common_frame(slide, "SOLUTION ILLUSTRATION")
    add_content_border(slide)

    add_box(slide, 0.92, 3.58, 1.55, 0.62, "MAIL / CLIENT INPUT", fill=PANEL, line=RGBColor(190, 205, 214), size=12)
    server = add_box(slide, 3.62, 2.36, 1.95, 2.85, content["project"]["solution_name"], fill=RGBColor(203, 216, 226), line=RGBColor(180, 195, 204), size=12)
    add_text(slide, 4.02, 2.05, 1.15, 0.2, "SERVER", size=11, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(slide, 3.85, 3.15, 1.48, 0.25, "WEB APP", size=11, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(slide, 3.85, 3.75, 1.48, 0.25, "BUSINESS ENGINE", size=10, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(slide, 3.85, 4.42, 1.48, 0.25, "DATABASE", size=11, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_arrow(slide, 2.48, 3.84, 3.58, 3.84, color=INK)
    add_arrow(slide, 3.58, 4.25, 2.48, 4.25, color=INK)

    modules = content["solution_modules"][:4]
    colors = [RED, BLUE, TEAL, BROWN]
    y_values = [1.85, 3.1, 4.28, 5.25]
    for index, module in enumerate(modules):
        y = y_values[index]
        color = colors[index % len(colors)]
        add_arrow(slide, 5.6, 3.6 + (index * 0.16), 8.08, y + 0.18, color=color)
        add_text(slide, 8.2, y, 3.0, 0.28, module.upper(), size=15, bold=True, color=color)
        add_box(slide, 10.65, y - 0.25, 0.82, 0.55, "APP", fill=RGBColor(238, 246, 242), line=color, size=10)


def slide_ui_mockups(prs: Presentation, content: dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_common_frame(slide, "UI MOCKUPS")
    add_content_border(slide)
    screens = content["ui_mockups"][:4]
    positions = [(1.15, 1.55), (6.7, 1.55), (1.15, 3.8), (6.7, 3.8)]
    for index, screen in enumerate(screens):
        x, y = positions[index]
        outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.75), Inches(1.75))
        set_fill(outer, RGBColor(246, 248, 248))
        set_line(outer, RGBColor(180, 190, 195), 1)
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(4.75), Inches(0.35))
        set_fill(header, GREEN)
        header.line.fill.background()
        add_text(slide, x + 0.12, y + 0.06, 4.45, 0.22, screen, size=13, bold=True, color=RGBColor(0, 0, 0))
        add_box(slide, x + 0.28, y + 0.62, 1.2, 0.3, "Field", fill=WHITE, line=RGBColor(210, 216, 216), size=8, bold=False)
        add_box(slide, x + 0.28, y + 1.05, 1.2, 0.3, "Field", fill=WHITE, line=RGBColor(210, 216, 216), size=8, bold=False)
        add_box(slide, x + 2.0, y + 0.62, 2.25, 0.75, "Data / Action Area", fill=RGBColor(232, 240, 247), line=BLUE, size=9, bold=True)


def slide_prerequisites(prs: Presentation, content: dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_common_frame(slide, "PREREQUISITES")
    add_content_border(slide)
    add_table(slide, 1.0, 1.55, 11.25, 4.45, ["Category", "Customer Responsibility", "Remarks"], content["prerequisites"], font_size=12)


def slide_delivery(prs: Presentation, content: dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_common_frame(slide, "DELIVERY SCHEDULE")
    add_content_border(slide)
    add_table(slide, 1.0, 1.55, 11.25, 4.45, ["Week", "Activity", "Owner"], content["delivery_schedule"], font_size=12)


def slide_contact(prs: Presentation, content: dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_common_frame(slide, "CONTACT MATRIX")
    add_content_border(slide)
    add_table(slide, 1.0, 1.55, 11.25, 4.45, ["Role", "Name", "Organization"], content["contact_matrix"], font_size=12)


def slide_thank_you(prs: Presentation, content: dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_common_frame(slide, "")
    add_text(slide, 1.1, 2.45, 11.1, 0.75, "Thank you", size=44, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(slide, 1.1, 3.35, 11.1, 0.35, "Ecosoft Zolutions Pvt Ltd", size=21, bold=True, color=MUTED, align=PP_ALIGN.CENTER)


def build_bom_pptx(content: dict[str, Any], output_dir: Path) -> tuple[Path, Path, Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs, content)
    slide_text(prs, "Requirements / Problem Statement", content["requirements"])
    slide_text(prs, "Proposed Solution", content["proposed_solution"])
    slide_technology(prs, content)
    slide_process_flow(prs, content)
    slide_solution_illustration(prs, content)
    slide_ui_mockups(prs, content)
    slide_text(prs, "Benefits", content["benefits"])
    slide_prerequisites(prs, content)
    slide_delivery(prs, content)
    slide_text(prs, "Customer Scope & Importants", content["customer_scope"])
    slide_contact(prs, content)
    slide_text(prs, "Business Model", content["business_model"])
    slide_text(prs, "Terms & Conditions", content["terms_conditions"])
    slide_thank_you(prs, content)

    project = content["project"]
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    base_name = safe_name(f'{project["project_code"]}_{project["customer_name"]}_{project["solution_name"]}_{stamp}')
    pptx_path = output_dir / f"{base_name}.pptx"
    json_path = output_dir / f"{base_name}.json"
    mmd_path = output_dir / f"{base_name}_process_flow.mmd"

    prs.save(pptx_path)
    json_path.write_text(json.dumps(content, indent=2), encoding="utf-8")
    mmd_path.write_text(content["process_mermaid"], encoding="utf-8")
    return pptx_path, json_path, mmd_path
